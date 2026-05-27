import argparse
import re
from pathlib import Path
from typing import List, Dict, Optional

from pptx import Presentation
from pptx.util import Inches
from PIL import Image, UnidentifiedImageError


SLIDE_HEADER_RE = re.compile(r"^##\s+Slide\s+\d+\s*-\s*(.+?)\s*\(([^)]+)\)\s*$")
IMAGE_TOKEN_RE = re.compile(r"([A-Za-z0-9_./\\-]+\.(?:png|jpg|jpeg|bmp|gif|tif|tiff))", re.IGNORECASE)
IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".bmp", ".gif", ".tif", ".tiff"}


def parse_slides_md(markdown_path: Path) -> List[Dict]:
    lines = markdown_path.read_text(encoding="utf-8").splitlines()

    slides = []
    current = None
    section = None

    def flush_current():
        if current:
            slides.append(current)

    for raw in lines:
        line = raw.rstrip()

        header_match = SLIDE_HEADER_RE.match(line)
        if header_match:
            flush_current()
            current = {
                "title": header_match.group(1).strip(),
                "time": header_match.group(2).strip(),
                "content": [],
                "visual": [],
                "notes": [],
            }
            section = None
            continue

        if current is None:
            continue

        if line.strip() == "**On-slide content**":
            section = "content"
            continue
        if line.strip() == "**Visual**":
            section = "visual"
            continue
        if line.strip() == "**Presenter notes**":
            section = "notes"
            continue
        if line.strip().startswith("---"):
            section = None
            continue

        stripped = line.strip()
        if not stripped:
            continue

        if section in {"content", "visual", "notes"}:
            # Accept bullets and numbered lines; keep readable text
            bullet = re.match(r"^[-*]\s+(.*)$", stripped)
            number = re.match(r"^\d+\.\s+(.*)$", stripped)
            nested = re.match(r"^-\s+(.*)$", stripped)

            text = None
            if bullet:
                text = bullet.group(1).strip()
            elif number:
                text = number.group(1).strip()
            elif nested:
                text = nested.group(1).strip()
            else:
                text = stripped

            if text:
                current[section].append(text)

    flush_current()
    return slides


def extract_image_paths(visual_lines: List[str]) -> List[str]:
    paths = []
    for line in visual_lines:
        # Prefer backtick paths from slides.md
        backticks = re.findall(r"`([^`]+)`", line)
        if backticks:
            for token in backticks:
                if Path(token).suffix.lower() in IMAGE_SUFFIXES:
                    paths.append(token)
            continue

        # If no backticks, pull image-like tokens only.
        for token in IMAGE_TOKEN_RE.findall(line):
            paths.append(token.strip())
    return paths


def resolve_image_path(raw_path: str, workspace_root: Path, thesis_root: Path) -> Optional[Path]:
    cleaned = raw_path.replace("\\", "/").strip()

    # Remove "Left:" etc.
    cleaned = re.sub(r"^[A-Za-z ]+:\s*", "", cleaned).strip()

    candidates = []

    def add_candidate(p: Path):
        if p not in candidates:
            candidates.append(p)

    add_candidate(workspace_root / cleaned)
    add_candidate(thesis_root / cleaned)

    if cleaned.startswith("thesis/"):
        stripped = cleaned[len("thesis/") :]
        add_candidate(workspace_root / stripped)
        add_candidate(thesis_root / stripped)

    # Auto-fallback: if source path points to figures/images, try out/figures/images mirror.
    normalized = cleaned
    if normalized.startswith("thesis/"):
        normalized = normalized[len("thesis/") :]

    figure_marker = "figures/images/"
    idx = normalized.find(figure_marker)
    if idx >= 0:
        rel_after_marker = normalized[idx + len(figure_marker) :]
        add_candidate(workspace_root / "out" / "figures" / "images" / rel_after_marker)

    # Also try relative-to-thesis and relative-to-workspace variants for convenience.
    add_candidate(thesis_root / cleaned)
    add_candidate(workspace_root / cleaned)

    for c in candidates:
        if c.exists() and c.is_file():
            return c
    return None


def is_valid_image(path: Path) -> bool:
    if path.suffix.lower() not in IMAGE_SUFFIXES:
        return False
    try:
        with Image.open(path) as img:
            img.verify()
        return True
    except (UnidentifiedImageError, OSError):
        return False


def add_images(slide, image_paths: List[Path], slide_label: str):
    # Keep image area in bottom half
    if not image_paths:
        return

    valid_paths = []
    for p in image_paths:
        if is_valid_image(p):
            valid_paths.append(p)
        else:
            print(f"WARN: {slide_label}: skipping invalid image file: {p}")

    if not valid_paths:
        return

    if len(valid_paths) == 1:
        slots = [(Inches(0.8), Inches(3.0), Inches(8.0))]
    elif len(valid_paths) == 2:
        slots = [
            (Inches(0.6), Inches(3.0), Inches(4.3)),
            (Inches(5.0), Inches(3.0), Inches(4.3)),
        ]
    else:
        slots = [
            (Inches(0.6), Inches(3.0), Inches(2.8)),
            (Inches(3.35), Inches(3.0), Inches(2.8)),
            (Inches(6.1), Inches(3.0), Inches(2.8)),
        ]

    for p, (left, top, width) in zip(valid_paths[:3], slots):
        try:
            slide.shapes.add_picture(str(p), left, top, width=width)
        except Exception as exc:
            print(f"WARN: {slide_label}: failed to insert image '{p}': {exc}")


def build_pptx(
    slides: List[Dict],
    output_path: Path,
    workspace_root: Path,
    thesis_root: Path,
    template_path: Optional[Path] = None,
):
    if template_path and template_path.exists():
        try:
            prs = Presentation(str(template_path))
        except Exception as exc:
            print(f"WARN: Failed to load template '{template_path}'. Using default template. Reason: {exc}")
            prs = Presentation()
    else:
        if template_path:
            print(f"WARN: Template not found: {template_path}. Using default template.")
        prs = Presentation()

    # Fallback layout indexes are usually stable:
    # 0 = title slide, 1 = title + content
    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]

    for idx, s in enumerate(slides, start=1):
        layout = title_layout if idx == 1 else content_layout
        slide = prs.slides.add_slide(layout)

        # Title
        slide.shapes.title.text = f"{s['title']} ({s['time']})"

        # Body for non-title slides
        body_placeholder = None
        if len(slide.placeholders) > 1:
            body_placeholder = slide.placeholders[1]

        if body_placeholder:
            tf = body_placeholder.text_frame
            tf.clear()

            if s["content"]:
                tf.text = s["content"][0]
                for item in s["content"][1:]:
                    p = tf.add_paragraph()
                    p.text = item
                    p.level = 0

        # Notes
        if s["notes"]:
            notes = slide.notes_slide.notes_text_frame
            notes.clear()
            notes.text = "\n".join(f"- {n}" for n in s["notes"])

        # Images
        slide_label = f"Slide {idx} ({s['title']})"
        raw_images = extract_image_paths(s["visual"])
        resolved = []
        for raw in raw_images:
            p = resolve_image_path(raw, workspace_root, thesis_root)
            if p:
                resolved.append(p)
            else:
                print(f"WARN: {slide_label}: could not resolve image path '{raw}'")

        # Keep max 3 images
        add_images(slide, resolved[:3], slide_label)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def main():
    parser = argparse.ArgumentParser(description="Generate dissertation slides from slides.md")
    parser.add_argument(
        "--input",
        default="slides.md",
        help="Path to slides markdown file (default: slides.md)",
    )
    parser.add_argument(
        "--output",
        default="MSA_Detector_Presentation_Draft.pptx",
        help="Output .pptx filename (default: MSA_Detector_Presentation_Draft.pptx)",
    )
    parser.add_argument(
        "--template",
        default=None,
        help="Optional .pptx template path (use your school template later)",
    )
    args = parser.parse_args()

    script_dir = Path(__file__).resolve().parent
    thesis_root = script_dir.parent
    workspace_root = thesis_root.parent

    input_path = (script_dir / args.input).resolve()
    output_path = (script_dir / args.output).resolve()
    template_path = Path(args.template).resolve() if args.template else None

    if template_path and not template_path.exists():
        print(f"WARN: Template not found: {template_path}. Using default template.")
        template_path = None

    slides = parse_slides_md(input_path)
    if not slides:
        raise ValueError(f"No slides parsed from: {input_path}")

    build_pptx(
        slides=slides,
        output_path=output_path,
        workspace_root=workspace_root,
        thesis_root=thesis_root,
        template_path=template_path,
    )
    print(f"Created: {output_path}")


if __name__ == "__main__":
    main()